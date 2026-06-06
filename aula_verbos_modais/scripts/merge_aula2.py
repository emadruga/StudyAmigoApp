import copy
from pptx import Presentation

def merge_pptx(files, output):
    merged = Presentation(files[0])
    for src_path in files[1:]:
        src = Presentation(src_path)
        for slide in src.slides:
            new_slide = merged.slides.add_slide(merged.slide_layouts[0])
            sp_tree = new_slide.shapes._spTree
            for ph in new_slide.placeholders:
                sp_tree.remove(ph.element)
            for shape in slide.shapes:
                sp_tree.append(copy.deepcopy(shape.element))
            try:
                src_bg = slide.background.fill._element
                dst_bg = new_slide.background.fill._element
                dst_bg.clear()
                for child in src_bg:
                    dst_bg.append(copy.deepcopy(child))
            except Exception:
                pass
    merged.save(output)
    print(f"✅ Saved: {output}  ({len(merged.slides)} slides)")

base = "/Users/emadruga/proj/StudyAmigoApp/aula_verbos_modais/"
merge_pptx(
    [base + "Aula2_PT_Slides_01-10.pptx",
     base + "Aula2_PT_Slides_11-20.pptx",
     base + "Aula2_PT_Slides_21-30.pptx"],
    base + "Aula2_COMPLETA.pptx"
)
