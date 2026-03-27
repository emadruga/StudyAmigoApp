"""
Gera os primeiros 11 slides de E02 em PPTX.
Uso: python gen_slides.py
Saída: E02_slides_1_11.pptx  (mesmo diretório)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Paleta de cores ──────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x1E, 0x1E, 0x2E)   # fundo escuro (quase preto-azul)
ACCENT     = RGBColor(0xF9, 0xC7, 0x42)   # amarelo-ouro (Brasil)
GREEN      = RGBColor(0x2E, 0xCC, 0x71)   # verde (Brasil / OK)
RED        = RGBColor(0xE7, 0x4C, 0x3C)   # vermelho (NÃO / atenção)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY   = RGBColor(0x44, 0x44, 0x55)   # células de tabela (fundo)
CODE_BG    = RGBColor(0x2A, 0x2A, 0x3E)   # fundo de bloco de código

# Dimensões widescreen
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Helpers ──────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    """Retorna um slide em branco com fundo escuro."""
    layout = prs.slide_layouts[6]          # layout completamente em branco
    slide  = prs.slides.add_slide(layout)
    fill   = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    return slide


def add_textbox(slide, text, x, y, w, h,
                font_size=24, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf    = txBox.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic    = italic
    return txBox


def add_rect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def slide_title_bar(slide, title_text, subtitle_text=None):
    """Faixa superior amarela com título."""
    bar_h = Inches(1.1)
    add_rect(slide, 0, 0, SLIDE_W, bar_h, ACCENT)
    add_textbox(slide, title_text,
                Inches(0.35), Inches(0.12), Inches(12.6), Inches(0.85),
                font_size=32, bold=True, color=DARK_BG, align=PP_ALIGN.LEFT)
    if subtitle_text:
        add_textbox(slide, subtitle_text,
                    Inches(0.35), Inches(1.15), Inches(12.6), Inches(0.45),
                    font_size=18, color=LIGHT_GRAY, italic=True)


def add_simple_table(slide, headers, rows,
                     x=Inches(0.4), y=Inches(1.65),
                     col_widths=None, row_height=Inches(0.52)):
    """Desenha uma tabela simples como retângulos + texto."""
    n_cols = len(headers)
    if col_widths is None:
        col_widths = [(SLIDE_W - Inches(0.8)) / n_cols] * n_cols

    HDR_BG   = RGBColor(0x2E, 0x46, 0x6E)
    ROW_BG1  = RGBColor(0x26, 0x26, 0x3A)
    ROW_BG2  = RGBColor(0x1E, 0x1E, 0x2E)
    BORDER   = RGBColor(0x44, 0x44, 0x66)

    # Cabeçalho
    cx = x
    for i, hdr in enumerate(headers):
        add_rect(slide, cx, y, col_widths[i], row_height, HDR_BG, BORDER)
        add_textbox(slide, hdr,
                    cx + Inches(0.07), y + Inches(0.07),
                    col_widths[i] - Inches(0.14), row_height - Inches(0.1),
                    font_size=16, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)
        cx += col_widths[i]

    # Linhas de dados
    for r_idx, row in enumerate(rows):
        cy = y + row_height * (r_idx + 1)
        bg = ROW_BG1 if r_idx % 2 == 0 else ROW_BG2
        cx = x
        for c_idx, cell in enumerate(row):
            add_rect(slide, cx, cy, col_widths[c_idx], row_height, bg, BORDER)
            # texto pode ser (text, color) ou plain text
            if isinstance(cell, tuple):
                cell_text, cell_color = cell
            else:
                cell_text, cell_color = cell, WHITE
            add_textbox(slide, cell_text,
                        cx + Inches(0.07), cy + Inches(0.06),
                        col_widths[c_idx] - Inches(0.14), row_height - Inches(0.08),
                        font_size=15, color=cell_color, align=PP_ALIGN.LEFT)
            cx += col_widths[c_idx]


def add_callout(slide, text, x, y, w, h, bg=MID_GRAY, icon="ℹ"):
    """Caixa de destaque / callout."""
    add_rect(slide, x, y, w, h, bg)
    add_textbox(slide, f"{icon}  {text}",
                x + Inches(0.15), y + Inches(0.1),
                w - Inches(0.25), h - Inches(0.15),
                font_size=17, color=WHITE, italic=True)


def add_code_block(slide, code_text, x, y, w, h, font_size=14):
    """Retângulo escuro com texto monoespaçado."""
    add_rect(slide, x, y, w, h, CODE_BG)
    txBox = slide.shapes.add_textbox(
        x + Inches(0.2), y + Inches(0.15),
        w - Inches(0.35), h - Inches(0.25))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for line in code_text.split("\n"):
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = line if line else " "
        run.font.size = Pt(font_size)
        run.font.color.rgb = RGBColor(0xA8, 0xFF, 0xC8)   # verde claro
        run.font.name = "Courier New"


# ── SLIDE 1 — Capa ───────────────────────────────────────────────────────────

def make_slide1(prs):
    slide = blank_slide(prs)

    # Banda decorativa esquerda
    add_rect(slide, 0, 0, Inches(0.35), SLIDE_H, ACCENT)

    # Título principal
    add_textbox(slide, "Exercício E02",
                Inches(0.65), Inches(1.6), Inches(11), Inches(1.2),
                font_size=54, bold=True, color=ACCENT)

    # Subtítulo
    add_textbox(slide, "Criando seus primeiros flashcards",
                Inches(0.65), Inches(2.85), Inches(11), Inches(0.8),
                font_size=32, color=WHITE)

    # Linha separadora
    add_rect(slide, Inches(0.65), Inches(3.75), Inches(9), Inches(0.04), ACCENT)

    # Detalhes
    add_textbox(slide, "StudyAmigo  ·  Abril 2026",
                Inches(0.65), Inches(3.95), Inches(8), Inches(0.5),
                font_size=20, color=LIGHT_GRAY)

    add_textbox(slide, "Biotecnologia  ·  Metrologia  ·  Segurança Cibernética",
                Inches(0.65), Inches(4.5), Inches(10), Inches(0.5),
                font_size=18, color=LIGHT_GRAY, italic=True)

    return slide


# ── SLIDE 2 — O que mudou de E01 para E02 ────────────────────────────────────

def make_slide2(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "O que mudou de E01 para E02")

    add_textbox(slide, "Em E01 você revisou cartões prontos.  Em E02 você vai criar os seus próprios.",
                Inches(0.4), Inches(1.2), Inches(12.5), Inches(0.5),
                font_size=19, color=LIGHT_GRAY, italic=True)

    headers = ["", "E01", "E02"]
    rows = [
        ["Baralho",  "Pré-carregado pelo professor", ("Você cria o baralho", GREEN)],
        ["Atividade","Apenas revisar",               ("Ler → selecionar → criar → revisar", ACCENT)],
        ["Escolha",  "Sem escolha",                  ("Você decide o que aprender", GREEN)],
    ]
    col_widths = [Inches(2.8), Inches(4.5), Inches(5.6)]
    add_simple_table(slide, headers, rows,
                     x=Inches(0.4), y=Inches(1.75),
                     col_widths=col_widths, row_height=Inches(0.72))

    add_callout(slide,
                "Criar um cartão é uma habilidade. Esta aula ensina como fazer certo desde o início.",
                Inches(0.4), Inches(5.8), Inches(12.5), Inches(0.7))

    return slide


# ── SLIDE 3 — Estrutura de E02 ───────────────────────────────────────────────

def make_slide3(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "Estrutura do E02: três componentes")

    headers = ["Componente", "O que é", "Peso"]
    rows = [
        [("A — Baralho Compartilhado", ACCENT),
         "Revisão do baralho curado (Verbal Tenses — Batch 2)", "~30%"],
        [("B — Texto Nivelado", ACCENT),
         "Leitura do texto do seu tier + criação de cartões", "~40%"],
        [("C — Livre Escolha", ACCENT),
         "Material à sua escolha + mínimo de 5 cartões", "~30%"],
    ]
    col_widths = [Inches(3.2), Inches(7.8), Inches(1.9)]
    add_simple_table(slide, headers, rows,
                     x=Inches(0.4), y=Inches(1.3),
                     col_widths=col_widths, row_height=Inches(0.88))

    add_callout(slide, "Esta aula foca nos Componentes B e C.",
                Inches(0.4), Inches(5.9), Inches(6.0), Inches(0.6),
                icon="→")

    return slide


# ── SLIDE 4 — Qual é o seu tier? ─────────────────────────────────────────────

def make_slide4(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "Qual é o seu tier?",
                    "Cada aluno recebe um texto diferente")

    headers = ["Tier", "Perfil", "Estimativa da turma"]
    rows = [
        [("Tier 1 — Foundation", ACCENT),
         "Inglês mínimo ou passivo; reconhece cognatos", "~55–60%"],
        [("Tier 2 — Developing", GREEN),
         "Inglês escolar; lê frases simples", "~33–37%"],
        [("Tier 3 — Expanding", RGBColor(0x5D, 0xAD, 0xE8)),
         "Inglês com exposição significativa; lê parágrafos", "~4–9%"],
    ]
    col_widths = [Inches(2.9), Inches(7.5), Inches(2.5)]
    add_simple_table(slide, headers, rows,
                     x=Inches(0.4), y=Inches(1.6),
                     col_widths=col_widths, row_height=Inches(0.88))

    add_callout(slide, "O professor já informou o seu tier. Se tiver dúvida, pergunte agora.",
                Inches(0.4), Inches(5.85), Inches(12.5), Inches(0.65))

    return slide


# ── SLIDE 5 — Qual texto você vai usar ───────────────────────────────────────

def make_slide5(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "Qual texto você vai usar",
                    "Component B: texto por tier")

    headers = ["Tier", "Texto", "Tamanho"]
    rows = [
        [("Tier 1", ACCENT),         "Estória Conhecida — versão simplificada", "~200 palavras"],
        [("Tier 2", GREEN),           "Estória Conhecida — versão mais rica",    "~300 palavras"],
        [("Tier 3", RGBColor(0x5D, 0xAD, 0xE8)),
         "Texto autêntico — IA / Demis Hassabis (60 Minutes / NYT)", "400–600 palavras"],
    ]
    col_widths = [Inches(1.8), Inches(8.5), Inches(2.6)]
    add_simple_table(slide, headers, rows,
                     x=Inches(0.4), y=Inches(1.6),
                     col_widths=col_widths, row_height=Inches(0.88))

    add_callout(slide,
                "Por que uma estória conhecida?  Você já conhece o enredo em português — "
                "sua atenção fica toda no vocabulário.",
                Inches(0.4), Inches(5.85), Inches(12.5), Inches(0.65),
                icon="💡")

    return slide


# ── SLIDE 6 — O que é um cartão bem feito? ───────────────────────────────────

def make_slide6(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "O que é um cartão bem feito?",
                    "Três regras para todos os tiers")

    rules = [
        ("Regra 1 — Cartão atômico",
         "Cada cartão testa exatamente uma coisa."),
        ("Regra 2 — Contexto mínimo",
         "A frente deve ter contexto suficiente para ter só uma resposta correta."),
        ("Regra 3 — Verso produtivo",
         "O verso deve ter mais do que a resposta — uma âncora de memória."),
    ]

    y = Inches(1.4)
    colors = [ACCENT, GREEN, RGBColor(0x5D, 0xAD, 0xE8)]
    for i, (label, desc) in enumerate(rules):
        box_y = y + i * Inches(1.7)
        # Número à esquerda
        add_rect(slide, Inches(0.4), box_y, Inches(0.55), Inches(1.35), colors[i])
        add_textbox(slide, str(i + 1),
                    Inches(0.4), box_y + Inches(0.28),
                    Inches(0.55), Inches(0.8),
                    font_size=32, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
        # Corpo
        add_rect(slide, Inches(0.98), box_y, Inches(11.9), Inches(1.35), MID_GRAY)
        add_textbox(slide, label,
                    Inches(1.12), box_y + Inches(0.1),
                    Inches(11.5), Inches(0.45),
                    font_size=19, bold=True, color=colors[i])
        add_textbox(slide, desc,
                    Inches(1.12), box_y + Inches(0.55),
                    Inches(11.5), Inches(0.65),
                    font_size=17, color=WHITE)

    return slide


# ── SLIDE 7 — O formato correto da frente ────────────────────────────────────

def make_slide7(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "O formato correto da frente",
                    "A palavra-alvo vai em [MAIÚSCULO]")

    code_bad = (
        "❌  Frente ambígua:\n"
        '   "Engineers analyze sensor data."\n'
        "   (O que o cartão está testando?  analyze?  sensor?  data?)"
    )
    code_good = (
        "✅  Frente correta:\n"
        '   "Engineers [ANALYZE] sensor data."\n'
        "   (Claro: o cartão testa a palavra 'analyze')"
    )

    add_code_block(slide, code_bad,  Inches(0.4), Inches(1.45), Inches(12.5), Inches(1.65))
    add_code_block(slide, code_good, Inches(0.4), Inches(3.25), Inches(12.5), Inches(1.65))

    add_callout(slide,
                "Importante: a frente é uma frase com contexto, não a palavra isolada.",
                Inches(0.4), Inches(5.9), Inches(12.5), Inches(0.6))

    return slide


# ── SLIDE 8 — O que NÃO fazer ────────────────────────────────────────────────

def make_slide8(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "O que NÃO fazer",
                    "Cartão antigo — todos os tiers devem evitar")

    code = (
        "Frente:  process\n"
        "Verso:   processo / processar"
    )
    add_code_block(slide, code, Inches(0.4), Inches(1.45), Inches(12.5), Inches(1.2))

    reasons = [
        ("Sem contexto",
         "Você reconhece a palavra, mas não aprende a usá-la."),
        ("Muito fácil",
         "Após 2–3 revisões você clica Easy — SM-2 agenda para 60 dias sem aprendizado real."),
        ("Não entra na memória de longo prazo",
         "Palavra isolada não cria conexões duradouras no cérebro."),
    ]

    y = Inches(2.85)
    for label, desc in reasons:
        add_rect(slide, Inches(0.4), y, Inches(0.08), Inches(0.55), RED)
        add_textbox(slide, f"{label}:  {desc}",
                    Inches(0.6), y,
                    Inches(12.2), Inches(0.55),
                    font_size=17, color=WHITE)
        y += Inches(0.68)

    add_callout(slide,
                "Por que é ruim? O SM-2 agenda como 'aprendido' algo que você apenas reconheceu.",
                Inches(0.4), Inches(5.9), Inches(12.5), Inches(0.6),
                bg=RGBColor(0x5E, 0x1A, 0x1A), icon="⚠")

    return slide


# ── SLIDE 9 — Tier 1: visão geral ────────────────────────────────────────────

def make_slide9(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "Tier 1 — Foundation: visão geral")

    headers = ["Elemento", "Instrução"]
    rows = [
        ["Texto",          "Estória Conhecida simplificada (~200 palavras)"],
        [("Lista de palavras", ACCENT),
         ("10 palavras fornecidas pelo professor", ACCENT)],
        ["Meta de cartões", "mín. 5  ·  alvo 7–8  ·  máx. 10"],
        ["Sessões",         "1–2 sessões ao longo do exercício"],
        ["Dicionário",      "Google Translate é suficiente"],
    ]
    col_widths = [Inches(2.8), Inches(10.1)]
    add_simple_table(slide, headers, rows,
                     x=Inches(0.4), y=Inches(1.3),
                     col_widths=col_widths, row_height=Inches(0.72))

    add_callout(slide,
                "Sua tarefa é aprender a criar bons cartões — não escolher quais palavras estudar. "
                "O professor já escolheu as palavras para você.",
                Inches(0.4), Inches(6.0), Inches(12.5), Inches(0.6),
                icon="→")

    return slide


# ── SLIDE 10 — Tier 1: passo a passo ─────────────────────────────────────────

def make_slide10(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "Tier 1 — Foundation: passo a passo",
                    "Como criar um cartão")

    steps = [
        "Pegue a lista de 10 palavras-alvo do professor.",
        "Escolha uma palavra da lista.",
        "Procure 1 frase exemplo com essa palavra (Google ou Cambridge).",
        'Troque uma palavra de contexto da frase por algo da sua vida:\n'
        '     "Scientists" → "Engineers" / "Biologists" / "I"',
        'Crie o cartão no StudyAmigo:\n'
        '     Frente: sua frase com a palavra em [MAIÚSCULO]\n'
        '     Verso: palavra = tradução em português',
        "Repita para mais 4–7 palavras da lista.",
        "Revise imediatamente após criar.",
    ]

    y = Inches(1.35)
    for i, step in enumerate(steps):
        step_h = Inches(0.62) if "\n" not in step else Inches(0.88)
        # Número
        add_rect(slide, Inches(0.4), y, Inches(0.45), step_h,
                 ACCENT if i < 3 else MID_GRAY)
        add_textbox(slide, str(i + 1),
                    Inches(0.4), y + Inches(0.08),
                    Inches(0.45), step_h - Inches(0.1),
                    font_size=17, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
        # Texto
        add_textbox(slide, step,
                    Inches(0.95), y + Inches(0.08),
                    Inches(12.0), step_h - Inches(0.1),
                    font_size=16, color=WHITE)
        y += step_h + Inches(0.06)

    return slide


# ── SLIDE 11 — Tier 1: exemplo completo ──────────────────────────────────────

def make_slide11(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "Tier 1 — Exemplo completo",
                    "Palavra: gather")

    code = (
        'Frase do dicionário:\n'
        '  "The girl gathered flowers in the forest."\n'
        '\n'
        'Substituição (aluno de Metrologia):\n'
        '  "I gathered data from the sensors in the lab."\n'
        '\n'
        'Cartão criado:\n'
        '  Frente:  "I [GATHERED] data from the sensors in the lab."\n'
        '  Verso:   gather / gathered  =  coletar, reunir'
    )
    add_code_block(slide, code, Inches(0.4), Inches(1.4), Inches(12.5), Inches(2.9),
                   font_size=15)

    bullets = [
        ("Você modificou a frase",       "→ seu cérebro processou ativamente"),
        ("A frase fala da sua vida",      "→ âncora de memória pessoal"),
        ("A palavra está marcada",        "→ SM-2 testa exatamente o que você aprendeu"),
    ]

    y = Inches(4.5)
    for label, effect in bullets:
        add_rect(slide, Inches(0.4), y, Inches(0.08), Inches(0.5), GREEN)
        add_textbox(slide, f"{label}  {effect}",
                    Inches(0.62), y + Inches(0.02),
                    Inches(12.3), Inches(0.5),
                    font_size=17, color=WHITE)
        y += Inches(0.6)

    return slide


# ── SLIDE 12 — Cambridge Learner's Dictionary ────────────────────────────────

def make_slide12(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "Ferramenta: Cambridge Learner's Dictionary",
                    "Por que o Cambridge, não o Google Translate?")

    headers = ["", "Google Translate", "Cambridge Learner's Dictionary"]
    rows = [
        ["O que fornece",
         "Tradução direta",
         ("Tradução + classe gramatical + colocações + exemplo", GREEN)],
        ["Ideal para",
         "Entender o significado rapidamente",
         ("Criar o verso completo do cartão", GREEN)],
        ["Tier recomendado",
         "Tier 1",
         ("Tier 2 e Tier 3", ACCENT)],
    ]
    col_widths = [Inches(2.4), Inches(4.5), Inches(6.0)]
    add_simple_table(slide, headers, rows,
                     x=Inches(0.4), y=Inches(1.55),
                     col_widths=col_widths, row_height=Inches(0.72))

    # Bullets "O que o Cambridge mostra"
    add_textbox(slide, "O que o Cambridge mostra que o Google Translate não mostra:",
                Inches(0.4), Inches(4.3), Inches(12.5), Inches(0.4),
                font_size=17, bold=True, color=ACCENT)
    items = [
        "Classe gramatical:  (v.)  (n.)  (adj.)",
        "Nível de dificuldade da palavra:  A1 → C2",
        "Colocações comuns:  gather data  /  gather evidence  /  gather information",
        "Frase-exemplo editada por lexicógrafos",
    ]
    y = Inches(4.75)
    for item in items:
        add_rect(slide, Inches(0.4), y + Inches(0.1), Inches(0.08), Inches(0.32), ACCENT)
        add_textbox(slide, item, Inches(0.58), y, Inches(12.2), Inches(0.48),
                    font_size=16, color=WHITE)
        y += Inches(0.48)

    add_callout(slide,
                "Acesso gratuito:  dictionary.cambridge.org  →  pesquise a palavra  →  'Learner's Dictionary'",
                Inches(0.4), Inches(6.92), Inches(12.5), Inches(0.55),
                icon="🌐")
    return slide


# ── SLIDE 13 — Tier 2: visão geral ───────────────────────────────────────────

def make_slide13(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "Tier 2 — Developing: visão geral")

    headers = ["Elemento", "Instrução"]
    rows = [
        ["Texto",           "Estória Conhecida — versão rica (~300 palavras)"],
        [("Lista de palavras", ACCENT),
         ("20 palavras fornecidas pelo professor", ACCENT)],
        ["Meta de cartões", "mín. 8  ·  alvo 10–12  ·  máx. 15"],
        [("Sessões", RGBColor(0x5D, 0xAD, 0xE8)),
         ("2–3 sessões  (não tudo no mesmo dia)", RGBColor(0x5D, 0xAD, 0xE8))],
        [("Dicionário", GREEN),
         ("Cambridge Learner's Dictionary", GREEN)],
    ]
    col_widths = [Inches(2.8), Inches(10.1)]
    add_simple_table(slide, headers, rows,
                     x=Inches(0.4), y=Inches(1.3),
                     col_widths=col_widths, row_height=Inches(0.72))

    add_callout(slide,
                "Você escolhe quais das 20 palavras vai transformar em cartão. "
                "Escolher quais palavras estudar já é um exercício valioso.",
                Inches(0.4), Inches(6.0), Inches(12.5), Inches(0.6),
                icon="→")
    return slide


# ── SLIDE 14 — Tier 2: formato do verso ──────────────────────────────────────

def make_slide14(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "Tier 2 — O formato do verso",
                    "O verso do Tier 2 exige mais do que o Tier 1")

    code = (
        "Tier 1 (simples):\n"
        "  gather / gathered  =  coletar, reunir\n"
        "\n"
        "Tier 2 (completo):\n"
        "  gather (v.)  =  coletar, reunir\n"
        "  ☞ gather data  /  gather evidence  /  gather information\n"
        "  ─────────────────────────────────────────────────────────\n"
        '  "Nos coletamos dados de calibracao dos equipamentos a cada sessao."'
    )
    add_code_block(slide, code, Inches(0.4), Inches(1.45), Inches(12.5), Inches(3.1),
                   font_size=15)

    reqs = [
        ('Classe gramatical', '(v.)  (n.)  (adj.)'),
        ('Colocação', '☞ gather data  /  gather evidence'),
        ('Tradução da frase', 'separada por linha horizontal'),
    ]
    add_textbox(slide, "O verso do Tier 2 deve ter sempre:",
                Inches(0.4), Inches(4.7), Inches(12.5), Inches(0.38),
                font_size=17, bold=True, color=ACCENT)
    y = Inches(5.1)
    for label, detail in reqs:
        add_rect(slide, Inches(0.4), y + Inches(0.08), Inches(0.08), Inches(0.35), GREEN)
        add_textbox(slide, f"{label}:  {detail}",
                    Inches(0.58), y, Inches(12.2), Inches(0.48),
                    font_size=16, color=WHITE)
        y += Inches(0.5)
    return slide


# ── SLIDE 15 — Tier 2: passo a passo ─────────────────────────────────────────

def make_slide15(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "Tier 2 — Developing: passo a passo",
                    "Como criar um cartão")

    steps = [
        "Pegue a lista de 20 palavras-alvo do professor.",
        "Escolha uma palavra — aquela que você não saberia usar numa frase agora.",
        'Abra o Cambridge Dictionary (learner\'s version):\n'
        '     anote a classe gramatical  ·  anote uma colocação comum (☞)',
        "Use a Tradução Reversa para criar a frase:\n"
        '     escreva em português  →  cole no Google Translate  →  edite pelo menos um elemento',
        "Monte o cartão:\n"
        "     Frente: sua frase com a palavra em [MAIÚSCULO]\n"
        "     Verso: palavra (classe) = tradução  ☞ colocação  +  tradução da frase",
        "Repita para 7–11 palavras, em sessões separadas.",
        "Revise após cada sessão.",
    ]

    y = Inches(1.35)
    for i, step in enumerate(steps):
        n_lines = step.count("\n")
        step_h = Inches(0.58) if n_lines == 0 else Inches(0.58 + n_lines * 0.38)
        add_rect(slide, Inches(0.4), y, Inches(0.45), step_h,
                 ACCENT if i < 3 else MID_GRAY)
        add_textbox(slide, str(i + 1),
                    Inches(0.4), y + Inches(0.08),
                    Inches(0.45), step_h - Inches(0.1),
                    font_size=17, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
        add_textbox(slide, step,
                    Inches(0.95), y + Inches(0.06),
                    Inches(12.0), step_h - Inches(0.08),
                    font_size=15, color=WHITE)
        y += step_h + Inches(0.05)
    return slide


# ── SLIDE 16 — Tier 2: exemplo completo ──────────────────────────────────────

def make_slide16(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "Tier 2 — Exemplo completo",
                    "Palavra: suspicious")

    code = (
        "Cambridge:\n"
        "  suspicious (adj.) = desconfiado, suspeito\n"
        "  ☞ suspicious of  /  feel suspicious  /  suspicious behavior\n"
        "\n"
        'Frase em portugues:  "O engenheiro ficou desconfiado do resultado do sensor."\n'
        "\n"
        'Google Translate:    "The engineer was suspicious of the sensor reading."\n'
        "\n"
        'Edicao (personalizacao):\n'
        '  "We were suspicious of the calibration results after the first test."\n'
        "\n"
        "Cartao criado:\n"
        '  Frente: "We were [SUSPICIOUS] of the calibration results after the first test."\n'
        "  Verso:  suspicious (adj.) = desconfiado  ☞ suspicious of / feel suspicious\n"
        "          ──────────────────────────────────────────────────────────────────\n"
        '          "Ficamos desconfiados dos resultados de calibracao apos o primeiro teste."'
    )
    add_code_block(slide, code, Inches(0.4), Inches(1.4), Inches(12.5), Inches(5.75),
                   font_size=13)
    return slide


# ── SLIDE 17 — Tier 3: visão geral ───────────────────────────────────────────

BLUE_T3 = RGBColor(0x5D, 0xAD, 0xE8)

def make_slide17(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "Tier 3 — Expanding: visão geral")

    headers = ["Elemento", "Instrução"]
    rows = [
        ["Texto",
         "Texto autêntico (60 Minutes / NYT — IA)  ·  400–600 palavras"],
        [("Lista de palavras", BLUE_T3),
         ("Seleção autônoma — você sublinha o que não conhece", BLUE_T3)],
        ["Meta de cartões", "mín. 10  ·  alvo 13–15  ·  máx. 18"],
        ["Sessões",          "2–3 sessões"],
        [("Dicionário", GREEN),
         ("Cambridge  +  Merriam-Webster", GREEN)],
    ]
    col_widths = [Inches(2.8), Inches(10.1)]
    add_simple_table(slide, headers, rows,
                     x=Inches(0.4), y=Inches(1.3),
                     col_widths=col_widths, row_height=Inches(0.72))

    add_callout(slide,
                "O verso é inteiramente em inglês — sem tradução da palavra, apenas definição em inglês.",
                Inches(0.4), Inches(6.0), Inches(12.5), Inches(0.6),
                icon="→")
    return slide


# ── SLIDE 18 — Tier 3: passo a passo ─────────────────────────────────────────

def make_slide18(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "Tier 3 — Expanding: passo a passo",
                    "Como criar um cartão")

    steps = [
        "Leia o texto uma vez completa sem parar para buscar palavras.",
        "Releia e sublinhe todas as palavras que você não conhece ou conhece só parcialmente.",
        "Para cada palavra sublinhada, abra o Cambridge ou Merriam-Webster:\n"
        "     copie a classe gramatical e a definição em inglês  ·  escolha uma frase exemplo do dicionário",
        "Monte a frente: frase do dicionário com a palavra em [MAIÚSCULO].",
        "Monte o verso:\n"
        "     Linha 1: palavra (classe) = definição em inglês\n"
        "     Linha 2: ☞ colocação 1  /  colocação 2\n"
        "     Linha 3: tradução em português da frase da frente",
        "Crie os cartões no baralho PASSAGE_E02_TIER3.",
        "Revise todos os dias até o fim do exercício.",
    ]

    y = Inches(1.35)
    for i, step in enumerate(steps):
        n_lines = step.count("\n")
        step_h = Inches(0.58) if n_lines == 0 else Inches(0.58 + n_lines * 0.36)
        add_rect(slide, Inches(0.4), y, Inches(0.45), step_h,
                 BLUE_T3 if i < 3 else MID_GRAY)
        add_textbox(slide, str(i + 1),
                    Inches(0.4), y + Inches(0.08),
                    Inches(0.45), step_h - Inches(0.1),
                    font_size=17, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
        add_textbox(slide, step,
                    Inches(0.95), y + Inches(0.06),
                    Inches(12.0), step_h - Inches(0.08),
                    font_size=15, color=WHITE)
        y += step_h + Inches(0.05)
    return slide


# ── SLIDE 19 — Tier 3: exemplo completo ──────────────────────────────────────

def make_slide19(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "Tier 3 — Exemplo completo",
                    "Palavra: drawn  (texto sobre Hassabis / DeepMind)")

    code = (
        'Trecho do texto:\n'
        '  "Hassabis was initially drawn to neuroscience because he believed\n'
        '   understanding the brain was the most reliable path to achieving\n'
        '   artificial general intelligence."\n'
        '\n'
        'Dicionario (Cambridge):\n'
        '  draw (v.) = to attract or interest someone\n'
        '  "She was drawn to the idea of working abroad."\n'
        '\n'
        'Cartao criado:\n'
        '  Frente: "She was [DRAWN] to the idea of working abroad."\n'
        '  Verso:  draw (v.) = to attract or interest someone\n'
        '          ☞ be drawn to  /  draw attention  /  draw interest\n'
        '          ──────────────────────────────────────────────────\n'
        '          "Ela se sentiu atraida pela ideia de trabalhar no exterior."'
    )
    add_code_block(slide, code, Inches(0.4), Inches(1.4), Inches(12.5), Inches(5.75),
                   font_size=14)
    return slide


# ── SLIDE 20 — Comparação: os três tiers ─────────────────────────────────────

def make_slide20(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "Comparação: o que cada tier faz",
                    "Resumo dos três tiers lado a lado")

    headers = ["Critério", "Tier 1", "Tier 2", "Tier 3"]
    rows = [
        ["Texto",
         "Estória ~200",
         "Estória ~300",
         "Autêntico 400–600"],
        ["Quem escolhe as palavras",
         ("Professor (10)", ACCENT),
         ("Professor (20)", GREEN),
         ("O aluno", BLUE_T3)],
        ["Meta de cartões",      "5–8",       "8–12",      "10–15"],
        ["Como cria a frase",
         "Substituição de 1 contexto",
         "Tradução Reversa",
         "Frase do dicionário"],
        ["Idioma do verso",
         "Português",
         "PT + EN",
         "EN + tradução da frase"],
    ]
    col_widths = [Inches(2.9), Inches(2.9), Inches(3.0), Inches(4.0)]
    add_simple_table(slide, headers, rows,
                     x=Inches(0.27), y=Inches(1.35),
                     col_widths=col_widths, row_height=Inches(0.78))
    return slide


# ── SLIDE 21 — Quantos cartões você precisa criar ────────────────────────────

def make_slide21(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "Quantos cartões você precisa criar",
                    "Metas de produção por tier — Component B")

    headers = ["Tier", "Mínimo", "Alvo", "Máximo"]
    rows = [
        [("Tier 1", ACCENT),      "5",  ("7–8",   ACCENT), "10"],
        [("Tier 2", GREEN),        "8",  ("10–12", GREEN),  "15"],
        [("Tier 3", BLUE_T3),     "10", ("13–15", BLUE_T3), "18"],
    ]
    col_widths = [Inches(3.5), Inches(2.5), Inches(3.5), Inches(3.0)]
    add_simple_table(slide, headers, rows,
                     x=Inches(0.65), y=Inches(1.7),
                     col_widths=col_widths, row_height=Inches(1.0))

    add_callout(slide,
                "Importante: não crie todos os cartões no mesmo dia.\n"
                "A nota considera em quantos dias diferentes você criou cartões.\n"
                "Criar em 3 dias distintos vale mais do que criar 15 cartões em 1 hora.",
                Inches(0.4), Inches(5.55), Inches(12.5), Inches(1.0),
                icon="⚠")
    return slide


# ── SLIDE 22 — Nome do baralho ────────────────────────────────────────────────

def make_slide22(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "Nome do baralho (crítico!)",
                    "Nomeie o baralho exatamente como indicado")

    headers = ["Seu tier", "Nome do baralho"]
    rows = [
        [("Tier 1", ACCENT),  ("PASSAGE_E02_TIER1", ACCENT)],
        [("Tier 2", GREEN),   ("PASSAGE_E02_TIER2", GREEN)],
        [("Tier 3", BLUE_T3), ("PASSAGE_E02_TIER3", BLUE_T3)],
    ]
    col_widths = [Inches(3.0), Inches(9.5)]
    add_simple_table(slide, headers, rows,
                     x=Inches(0.4), y=Inches(1.5),
                     col_widths=col_widths, row_height=Inches(0.82))

    add_callout(slide,
                "⚠  O sistema usa o nome do baralho para calcular a nota. "
                "Um nome errado faz seus cartões serem contabilizados como Component C em vez de B.",
                Inches(0.4), Inches(4.05), Inches(12.5), Inches(0.75),
                bg=RGBColor(0x5E, 0x1A, 0x1A), icon="")

    add_textbox(slide, "Como criar o baralho no StudyAmigo:",
                Inches(0.4), Inches(5.0), Inches(12.5), Inches(0.4),
                font_size=17, bold=True, color=ACCENT)

    steps_text = (
        "1. Acesse StudyAmigo e faça login\n"
        "2. Na tela de baralhos, clique em  Novo Baralho\n"
        "3. Digite o nome exato da tabela acima\n"
        "4. Use este baralho para todos os cartões do texto de E02"
    )
    add_textbox(slide, steps_text,
                Inches(0.55), Inches(5.45), Inches(12.2), Inches(1.5),
                font_size=16, color=WHITE)
    return slide


# ── SLIDE 23 — Component C: material de livre escolha ────────────────────────

def make_slide23(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "Component C: material de livre escolha",
                    "Você também precisa de 5 cartões de material próprio")

    # Regras
    rules = [
        "Mínimo de 5 cartões",
        "O material deve ser um texto em inglês (não lista de palavras)",
        "Submeta a fonte junto com os cartões (URL, título, nome da música)",
    ]
    y = Inches(1.35)
    for rule in rules:
        add_rect(slide, Inches(0.4), y + Inches(0.1), Inches(0.08), Inches(0.32), ACCENT)
        add_textbox(slide, rule, Inches(0.58), y, Inches(12.2), Inches(0.48),
                    font_size=17, color=WHITE)
        y += Inches(0.5)

    # Tabela de materiais
    headers = ["Material", "Aceito?"]
    rows = [
        ["Letra de música",                              ("✅  (prefira vocabulário claro)", GREEN)],
        ["Transcrição de vídeo (YouTube, TED)",          ("✅  Recomendado", GREEN)],
        ["Artigo simplificado (BBC Learning, VOA)",      ("✅  Recomendado", GREEN)],
        ["Tradução de texto em português",               ("❌  Não aceito", RED)],
    ]
    col_widths = [Inches(8.5), Inches(4.4)]
    add_simple_table(slide, headers, rows,
                     x=Inches(0.4), y=Inches(3.05),
                     col_widths=col_widths, row_height=Inches(0.62))

    add_callout(slide,
                "O nome do baralho de Component C pode ser qualquer nome — só o Component B tem nome fixo.",
                Inches(0.4), Inches(6.28), Inches(12.5), Inches(0.58),
                icon="→")
    return slide


# ── SLIDE 24 — Sugestões de fontes para Component C ──────────────────────────

def make_slide24(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "Sugestões de fontes para Component C",
                    "Por onde começar?")

    tiers = [
        ("Tier 1", ACCENT, [
            "VOA Learning English — Beginners:  learningenglish.voanews.com",
            "Letras de músicas pop/folk com vocabulário simples",
            "Legendas de desenhos animados americanos",
        ]),
        ("Tier 2", GREEN, [
            "BBC Learning English — News Report:  bbc.co.uk/learningenglish",
            "Transcrições de vídeos TED-Ed",
            "Artigos de introdução da Wikipedia em inglês",
        ]),
        ("Tier 3", BLUE_T3, [
            "The Guardian  /  Scientific American  /  MIT Technology Review",
            "Transcrições de podcasts (Radiolab, Stuff You Should Know)",
            "Livros de não-ficção populares",
        ]),
    ]

    y = Inches(1.35)
    for tier_label, color, items in tiers:
        add_rect(slide, Inches(0.4), y, Inches(1.5), Inches(0.4), color)
        add_textbox(slide, tier_label,
                    Inches(0.4), y + Inches(0.03),
                    Inches(1.5), Inches(0.35),
                    font_size=16, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
        iy = y
        for item in items:
            add_rect(slide, Inches(2.05), iy + Inches(0.08), Inches(0.07), Inches(0.28), color)
            add_textbox(slide, item, Inches(2.2), iy, Inches(10.7), Inches(0.42),
                        font_size=15, color=WHITE)
            iy += Inches(0.43)
        y = iy + Inches(0.25)
    return slide


# ── SLIDE 25 — Erros mais comuns ──────────────────────────────────────────────

def make_slide25(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "Erros mais comuns — evite!",
                    "O que NÃO fazer")

    headers = ["❌  Erro", "✅  Correto"]
    rows = [
        ["Frente: palavra isolada  (process)",
         "Frente: frase com  [PROCESS]"],
        ["Verso: só tradução  (processo)",
         "Verso: tradução + exemplo / colocação"],
        ["Cartão de cognato óbvio  (technology = tecnologia)",
         "Pule cognatos que você já sabe"],
        ["Criar todos os cartões em 1 dia",
         "Distribuir em 2–3 sessões diferentes"],
        ["Nome do baralho errado  (Exercicio2)",
         "Nome exato:  PASSAGE_E02_TIER2"],
        ["Frente sem marcação da palavra-alvo",
         "Frente sempre com  [PALAVRA_EM_MAIÚSCULO]"],
    ]
    col_widths = [Inches(6.4), Inches(6.4)]
    add_simple_table(slide, headers, rows,
                     x=Inches(0.27), y=Inches(1.3),
                     col_widths=col_widths, row_height=Inches(0.65))
    return slide


# ── SLIDE 26 — O que acontece depois de criar ────────────────────────────────

def make_slide26(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "O que acontece depois de criar",
                    "O ciclo completo")

    # Fluxo como caixas encadeadas
    steps = [
        ("Criar cartão",              "Revisar imediatamente\n(1ª sessão de aprendizado)", ACCENT),
        ("SM-2 agenda",               "próxima revisão em 1–3 dias", MID_GRAY),
        ("Você volta e revisa",       "SM-2 ajusta o intervalo", GREEN),
        ("Cartão maduro",             "intervalo ≥ 21 dias\n(sua nota agradece)", RGBColor(0x5D, 0xAD, 0xE8)),
    ]

    box_w = Inches(2.8)
    box_h = Inches(1.7)
    gap   = Inches(0.45)
    start_x = Inches(0.55)
    start_y = Inches(1.7)

    for i, (title, desc, color) in enumerate(steps):
        bx = start_x + i * (box_w + gap)
        add_rect(slide, bx, start_y, box_w, box_h, color)
        add_textbox(slide, title,
                    bx + Inches(0.1), start_y + Inches(0.1),
                    box_w - Inches(0.2), Inches(0.5),
                    font_size=17, bold=True,
                    color=DARK_BG if color == ACCENT else WHITE,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, desc,
                    bx + Inches(0.1), start_y + Inches(0.65),
                    box_w - Inches(0.2), Inches(0.95),
                    font_size=14,
                    color=DARK_BG if color == ACCENT else WHITE,
                    align=PP_ALIGN.CENTER)
        # seta entre caixas
        if i < len(steps) - 1:
            ax = bx + box_w + Inches(0.05)
            add_textbox(slide, "→",
                        ax, start_y + Inches(0.6),
                        gap - Inches(0.05), Inches(0.55),
                        font_size=26, bold=True, color=LIGHT_GRAY,
                        align=PP_ALIGN.CENTER)

    add_callout(slide,
                "Revisar imediatamente após criar é obrigatório.\n"
                "Cartões criados e nunca revisados não entram no SM-2 e não geram nota.",
                Inches(0.4), Inches(5.8), Inches(12.5), Inches(0.75),
                bg=RGBColor(0x5E, 0x1A, 0x1A), icon="⚠")
    return slide


# ── SLIDE 27 — Checklist final ────────────────────────────────────────────────

def make_slide27(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "Checklist final",
                    "Antes de terminar o exercício, verifique:")

    # Dois painéis lado a lado
    def painel(title, items, color, x, y, w):
        add_rect(slide, x, y, w, Inches(0.42), color)
        add_textbox(slide, title,
                    x + Inches(0.1), y + Inches(0.05),
                    w - Inches(0.2), Inches(0.35),
                    font_size=17, bold=True, color=DARK_BG)
        iy = y + Inches(0.5)
        for item in items:
            add_rect(slide, x + Inches(0.1), iy + Inches(0.08),
                     Inches(0.22), Inches(0.22),
                     RGBColor(0x33, 0x33, 0x44))
            add_textbox(slide, item,
                        x + Inches(0.45), iy,
                        w - Inches(0.55), Inches(0.42),
                        font_size=15, color=WHITE)
            iy += Inches(0.46)

    comp_b = [
        "Criei o baralho com o nome exato do meu tier",
        "Criei pelo menos o mínimo de cartões do meu tier",
        "As frentes têm [PALAVRA EM MAIÚSCULO]",
        "Tier 2+: versos têm classe gramatical e ☞ colocação",
        "Criei cartões em mais de 1 dia (distribuição)",
        "Revisei os cartões logo após criar",
    ]
    comp_c = [
        "Criei pelo menos 5 cartões de material próprio",
        "Anotei a fonte do material (URL ou título)",
        "O material é um texto em inglês (não uma tradução)",
    ]

    painel("Component B", comp_b, ACCENT,
           Inches(0.4), Inches(1.3), Inches(6.3))
    painel("Component C", comp_c, GREEN,
           Inches(6.9), Inches(1.3), Inches(6.0))
    return slide


# ── SLIDE 28 — Resumo em uma frase por tier ───────────────────────────────────

def make_slide28(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "Resumo em uma frase por tier",
                    "O que levar para casa")

    blocks = [
        ("Tier 1", ACCENT,
         "Pegue a lista do professor, troque uma palavra de contexto na frase do dicionário, "
         "coloque a palavra-alvo em [MAIÚSCULO] na frente."),
        ("Tier 2", GREEN,
         "Use o Cambridge, escreva em português, traduza, edite uma coisa, "
         "coloque ☞ colocação no verso."),
        ("Tier 3", BLUE_T3,
         "Leia completo, sublinhe o que não sabe, use frase do dicionário, "
         "defina em inglês, ☞ colocação, tradução da frase no verso."),
    ]

    y = Inches(1.55)
    for tier_label, color, text in blocks:
        add_rect(slide, Inches(0.4), y, Inches(1.5), Inches(1.45), color)
        add_textbox(slide, tier_label,
                    Inches(0.4), y + Inches(0.45),
                    Inches(1.5), Inches(0.6),
                    font_size=18, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
        add_rect(slide, Inches(1.95), y, Inches(11.0), Inches(1.45), MID_GRAY)
        add_textbox(slide, text,
                    Inches(2.1), y + Inches(0.2),
                    Inches(10.7), Inches(1.1),
                    font_size=17, color=WHITE)
        y += Inches(1.65)
    return slide


# ── SLIDE 29 — Dúvidas? ───────────────────────────────────────────────────────

def make_slide29(prs):
    slide = blank_slide(prs)

    # Banda lateral
    add_rect(slide, 0, 0, Inches(0.35), SLIDE_H, ACCENT)

    add_textbox(slide, "Dúvidas?",
                Inches(0.65), Inches(1.3), Inches(11), Inches(1.1),
                font_size=54, bold=True, color=ACCENT)

    contacts = [
        ("Dúvidas sobre o tier",              "fale com o professor"),
        ("Dúvidas sobre o StudyAmigo",         "login, criar baralho → fale com o professor"),
        ("Dúvidas sobre o método de cartão",   "este documento + exemplos dos slides anteriores"),
    ]
    y = Inches(2.6)
    for label, detail in contacts:
        add_rect(slide, Inches(0.65), y + Inches(0.1), Inches(0.08), Inches(0.32), ACCENT)
        add_textbox(slide, f"{label}:  {detail}",
                    Inches(0.85), y,
                    Inches(11.8), Inches(0.5),
                    font_size=18, color=WHITE)
        y += Inches(0.62)

    add_rect(slide, Inches(0.65), Inches(4.45), Inches(9.0), Inches(0.04), MID_GRAY)

    add_textbox(slide, "Prazo de E02:",
                Inches(0.65), Inches(4.65), Inches(3.0), Inches(0.5),
                font_size=20, bold=True, color=ACCENT)
    add_textbox(slide, "Abril 2026  (data exata confirmada pelo professor)",
                Inches(3.7), Inches(4.65), Inches(8.5), Inches(0.5),
                font_size=20, color=WHITE)

    add_textbox(slide, "Documento gerado em Março 2026 · StudyAmigo / E02",
                Inches(0.65), Inches(6.8), Inches(10.0), Inches(0.4),
                font_size=13, color=LIGHT_GRAY, italic=True)
    return slide


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    makers = [
        make_slide1,  make_slide2,  make_slide3,  make_slide4,  make_slide5,
        make_slide6,  make_slide7,  make_slide8,  make_slide9,  make_slide10,
        make_slide11, make_slide12, make_slide13, make_slide14, make_slide15,
        make_slide16, make_slide17, make_slide18, make_slide19, make_slide20,
        make_slide21, make_slide22, make_slide23, make_slide24, make_slide25,
        make_slide26, make_slide27, make_slide28, make_slide29,
    ]

    for i, maker in enumerate(makers, start=1):
        print(f"  Gerando slide {i:02d}…")
        maker(prs)

    out = "placement_exam/planning_E02/E02_slides_completo.pptx"
    prs.save(out)
    print(f"\n✅  Salvo em: {out}  ({len(makers)} slides)")


if __name__ == "__main__":
    main()
