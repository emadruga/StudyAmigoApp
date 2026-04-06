"""
Gera os slides de Duas Passagens de Leitura — Tier 1 (E02) em PPTX.
Uso: python gen_slides_tier1_duas_passagens.py
Saída: E02_slides_Tier1_Duas_Passagem_v1.0.pptx  (mesmo diretório)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Paleta de cores (idêntica ao gen_slides.py) ──────────────────────────────
DARK_BG    = RGBColor(0x1E, 0x1E, 0x2E)
ACCENT     = RGBColor(0xF9, 0xC7, 0x42)
GREEN      = RGBColor(0x2E, 0xCC, 0x71)
RED        = RGBColor(0xE7, 0x4C, 0x3C)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY   = RGBColor(0x44, 0x44, 0x55)
CODE_BG    = RGBColor(0x2A, 0x2A, 0x3E)
BLUE       = RGBColor(0x5D, 0xAD, 0xE8)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Helpers (idênticos ao gen_slides.py) ────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    fill   = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    return slide


def add_textbox(slide, text, x, y, w, h,
                font_size=24, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf    = txBox.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size      = Pt(font_size)
    run.font.bold      = bold
    run.font.color.rgb = color
    run.font.italic    = italic
    return txBox


def add_rect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def slide_title_bar(slide, title_text, subtitle_text=None):
    bar_h = Inches(1.1)
    add_rect(slide, 0, 0, SLIDE_W, bar_h, ACCENT)
    add_textbox(slide, title_text,
                Inches(0.35), Inches(0.12), Inches(12.6), Inches(0.85),
                font_size=32, bold=True, color=DARK_BG, align=PP_ALIGN.LEFT)
    if subtitle_text:
        add_textbox(slide, subtitle_text,
                    Inches(0.35), Inches(1.15), Inches(12.6), Inches(0.45),
                    font_size=18, color=LIGHT_GRAY, italic=True)


def add_simple_table(slide, headers, rows,
                     x=Inches(0.4), y=Inches(1.65),
                     col_widths=None, row_height=Inches(0.52)):
    n_cols = len(headers)
    if col_widths is None:
        col_widths = [(SLIDE_W - Inches(0.8)) / n_cols] * n_cols

    HDR_BG  = RGBColor(0x2E, 0x46, 0x6E)
    ROW_BG1 = RGBColor(0x26, 0x26, 0x3A)
    ROW_BG2 = RGBColor(0x1E, 0x1E, 0x2E)
    BORDER  = RGBColor(0x44, 0x44, 0x66)

    cx = x
    for i, hdr in enumerate(headers):
        add_rect(slide, cx, y, col_widths[i], row_height, HDR_BG, BORDER)
        add_textbox(slide, hdr,
                    cx + Inches(0.07), y + Inches(0.07),
                    col_widths[i] - Inches(0.14), row_height - Inches(0.1),
                    font_size=16, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)
        cx += col_widths[i]

    for r_idx, row in enumerate(rows):
        cy = y + row_height * (r_idx + 1)
        bg = ROW_BG1 if r_idx % 2 == 0 else ROW_BG2
        cx = x
        for c_idx, cell in enumerate(row):
            add_rect(slide, cx, cy, col_widths[c_idx], row_height, bg, BORDER)
            if isinstance(cell, tuple):
                cell_text, cell_color = cell
            else:
                cell_text, cell_color = cell, WHITE
            add_textbox(slide, cell_text,
                        cx + Inches(0.07), cy + Inches(0.06),
                        col_widths[c_idx] - Inches(0.14), row_height - Inches(0.08),
                        font_size=15, color=cell_color, align=PP_ALIGN.LEFT)
            cx += col_widths[c_idx]


def add_callout(slide, text, x, y, w, h, bg=MID_GRAY, icon="ℹ"):
    add_rect(slide, x, y, w, h, bg)
    add_textbox(slide, f"{icon}  {text}",
                x + Inches(0.15), y + Inches(0.1),
                w - Inches(0.25), h - Inches(0.15),
                font_size=17, color=WHITE, italic=True)


def add_code_block(slide, code_text, x, y, w, h, font_size=14):
    add_rect(slide, x, y, w, h, CODE_BG)
    txBox = slide.shapes.add_textbox(
        x + Inches(0.2), y + Inches(0.15),
        w - Inches(0.35), h - Inches(0.25))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for line in code_text.split("\n"):
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = line if line else " "
        run.font.size = Pt(font_size)
        run.font.color.rgb = RGBColor(0xA8, 0xFF, 0xC8)
        run.font.name = "Courier New"


def section_divider(prs, title, subtitle=""):
    """Slide de transição entre seções."""
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, Inches(0.5), SLIDE_H, ACCENT)
    add_textbox(slide, title,
                Inches(0.85), Inches(2.4), Inches(11.5), Inches(1.4),
                font_size=44, bold=True, color=ACCENT)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.85), Inches(3.9), Inches(11.5), Inches(0.7),
                    font_size=24, color=LIGHT_GRAY, italic=True)
    return slide


# ── SLIDE 1 — Capa ───────────────────────────────────────────────────────────

def make_slide1(prs):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, Inches(0.35), SLIDE_H, ACCENT)
    add_textbox(slide, "Tier 1 — Duas Passagens de Leitura",
                Inches(0.65), Inches(1.4), Inches(11.5), Inches(1.3),
                font_size=46, bold=True, color=ACCENT)
    add_textbox(slide, "Como ler um texto em inglês sem travar e sem o Google Translate",
                Inches(0.65), Inches(2.85), Inches(11.5), Inches(0.8),
                font_size=28, color=WHITE)
    add_rect(slide, Inches(0.65), Inches(3.78), Inches(9), Inches(0.04), ACCENT)
    add_textbox(slide, "StudyAmigo  ·  E02  ·  Abril 2026",
                Inches(0.65), Inches(3.98), Inches(8), Inches(0.5),
                font_size=20, color=LIGHT_GRAY)
    add_textbox(slide, "Texto: A Seleção Brasileira (Versão 1 — História e Conquistas)",
                Inches(0.65), Inches(4.55), Inches(10.5), Inches(0.5),
                font_size=18, color=LIGHT_GRAY, italic=True)
    return slide


# ── SLIDE 2 — Visão geral das duas passagens ─────────────────────────────────

def make_slide2(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "As duas passagens de leitura — visão geral")

    headers = ["", "1ª Passagem", "2ª Passagem"]
    rows = [
        ["Objetivo",
         "Entender do que o texto trata",
         "Examinar as palavras-alvo nas suas frases"],
        ["Suporte",
         ("Nenhum — leitura autônoma", GREEN),
         "Lista do professor + dicionário (palavra a palavra)"],
        ["Foco",
         "Texto inteiro",
         "Somente as frases das 10 palavras-alvo"],
        ["Ferramenta",
         "Cognatos + contexto",
         ("Verbo → Sujeito → Objeto → Dicionário", ACCENT)],
        ["Produto",
         "Modelo mental do texto",
         ("10 cartões prontos para revisão", GREEN)],
        ["Tempo estimado",
         "5–8 minutos",
         "20–30 minutos"],
    ]
    col_widths = [Inches(2.2), Inches(4.8), Inches(5.9)]
    add_simple_table(slide, headers, rows,
                     x=Inches(0.4), y=Inches(1.15),
                     col_widths=col_widths, row_height=Inches(0.72))

    return slide


# ── SLIDE 3 — O problema: dois impulsos autodestrutivos ──────────────────────

def make_slide3(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "O problema: dois impulsos que travam o aluno Tier 1")

    # Impulso 1
    add_rect(slide, Inches(0.4), Inches(1.6), Inches(5.9), Inches(2.2),
             RGBColor(0x3A, 0x1A, 0x1A))
    add_textbox(slide, "① Paralisia",
                Inches(0.6), Inches(1.7), Inches(5.5), Inches(0.5),
                font_size=22, bold=True, color=RED)
    add_textbox(slide, "Para na primeira palavra desconhecida\ne não avança.",
                Inches(0.6), Inches(2.25), Inches(5.5), Inches(1.3),
                font_size=18, color=WHITE)

    # Impulso 2
    add_rect(slide, Inches(7.0), Inches(1.6), Inches(5.9), Inches(2.2),
             RGBColor(0x3A, 0x1A, 0x1A))
    add_textbox(slide, "② Atalho",
                Inches(7.2), Inches(1.7), Inches(5.5), Inches(0.5),
                font_size=22, bold=True, color=RED)
    add_textbox(slide, "Copia o texto inteiro e despeja\nno Google Translate ou ChatGPT.",
                Inches(7.2), Inches(2.25), Inches(5.5), Inches(1.3),
                font_size=18, color=WHITE)

    # Causa comum
    add_callout(slide,
                "Ambos vêm do mesmo medo: o medo de não entender nada.",
                Inches(0.4), Inches(4.05), Inches(12.5), Inches(0.65),
                bg=RGBColor(0x5A, 0x2A, 0x2A), icon="⚠")

    add_textbox(slide,
                "Esse medo existe porque o aluno acredita que 'entender' = traduzir cada palavra.",
                Inches(0.4), Inches(4.85), Inches(12.5), Inches(0.5),
                font_size=18, color=LIGHT_GRAY, italic=True)

    add_callout(slide,
                "A 1ª passagem precisa redefinir o que significa 'entender' um texto.",
                Inches(0.4), Inches(5.55), Inches(12.5), Inches(0.65),
                bg=RGBColor(0x1A, 0x3A, 0x2A), icon="→")

    return slide


# ── SLIDE 4 — Critério de sucesso da 1ª passagem ────────────────────────────

def make_slide4(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "O que conta como sucesso na 1ª passagem")

    add_callout(slide,
                "Critério de sucesso: ao final da leitura, o aluno consegue dizer "
                "em UMA FRASE em português o que o texto está falando.",
                Inches(0.4), Inches(1.3), Inches(12.5), Inches(0.85),
                bg=RGBColor(0x1A, 0x3A, 0x2A), icon="✓")

    add_textbox(slide, "Isso não é tradução. É compreensão global.",
                Inches(0.4), Inches(2.35), Inches(12.5), Inches(0.5),
                font_size=20, bold=True, color=ACCENT)

    # Instrução correta
    add_rect(slide, Inches(0.4), Inches(2.95), Inches(12.5), Inches(0.04), ACCENT)
    add_textbox(slide, "A instrução correta para o aluno:",
                Inches(0.4), Inches(3.1), Inches(12.5), Inches(0.45),
                font_size=18, bold=True, color=LIGHT_GRAY)

    add_code_block(slide,
                   '"Leia o texto uma vez do início ao fim. Não pare nas palavras que não conhece.\n'
                   'Não abra o Google Translate. Quando terminar, feche o texto e escreva uma\n'
                   'frase em português dizendo o que o texto está falando.\n'
                   'Se você conseguir fazer isso, a primeira leitura foi bem-feita."',
                   Inches(0.4), Inches(3.65), Inches(12.5), Inches(1.55), font_size=16)

    add_callout(slide,
                "Essa instrução dá ao aluno um critério concreto que ele consegue atingir sem ajuda externa.",
                Inches(0.4), Inches(5.45), Inches(12.5), Inches(0.65))

    return slide


# ── SLIDE 5 — Texto de referência ───────────────────────────────────────────

def make_slide5(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "Texto de referência — Tier 1, Versão 1",
                    "Primeiras 4 frases usadas na demonstração")

    add_code_block(slide,
                   'Brazil\'s national football team is called A Seleção.\n'
                   'The team played its first international game in 1914.\n'
                   'Since then, Brazil has competed in every single World Cup.\n'
                   'No other team has done this.',
                   Inches(0.4), Inches(1.65), Inches(12.5), Inches(2.1), font_size=18)

    add_textbox(slide, "Por que este texto funciona para Tier 1:",
                Inches(0.4), Inches(3.95), Inches(12.5), Inches(0.45),
                font_size=18, bold=True, color=ACCENT)

    items = [
        ("① Assunto familiar", "O aluno já conhece a Seleção — carga cognitiva de conteúdo = zero"),
        ("② Cognatos densos", "Brazil, national, international, football, World Cup"),
        ("③ Estrutura simples", "Frases curtas; uma oração principal por frase"),
        ("④ Datas e fatos", "1914, Copa do Mundo — pontos de apoio concretos"),
    ]
    y_start = Inches(4.5)
    for i, (label, desc) in enumerate(items):
        y = y_start + Inches(0.48) * i
        add_textbox(slide, label,
                    Inches(0.5), y, Inches(2.3), Inches(0.42),
                    font_size=16, bold=True, color=GREEN)
        add_textbox(slide, desc,
                    Inches(2.85), y, Inches(10.0), Inches(0.42),
                    font_size=16, color=WHITE)

    return slide


# ── SLIDE 6 — Frase 1: Brazil's national football team ──────────────────────

def make_slide6(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "1ª Passagem — Frase 1",
                    'Brazil\'s national football team is called A Seleção.')

    headers = ["Palavra / trecho", "O que o aluno reconhece"]
    rows = [
        [("Brazil", ACCENT),         "Brasil — sei"],
        ["national",                  "parece nacional — deve ser isso"],
        ["football team",             "time de futebol — sei"],
        ["is called",                 ("não sei, mas...", LIGHT_GRAY)],
        [("A Seleção", ACCENT),       "A Seleção Brasileira — reconhece imediatamente"],
    ]
    col_widths = [Inches(3.5), Inches(9.0)]
    add_simple_table(slide, headers, rows,
                     x=Inches(0.4), y=Inches(1.25),
                     col_widths=col_widths, row_height=Inches(0.6))

    add_callout(slide,
                "O aluno entende: o texto é sobre a Seleção Brasileira. "
                "Reconheceu o assunto antes de terminar a frase.",
                Inches(0.4), Inches(5.15), Inches(12.5), Inches(0.65),
                bg=RGBColor(0x1A, 0x3A, 0x2A), icon="✓")

    add_textbox(slide,
                "Nota pedagógica: o aluno não sabe o que is called significa — e não precisa saber ainda. "
                "O nome A Seleção no final entrega o sentido completo.",
                Inches(0.4), Inches(6.0), Inches(12.5), Inches(0.6),
                font_size=15, color=LIGHT_GRAY, italic=True)

    return slide


# ── SLIDE 7 — Frase 2: The team played ──────────────────────────────────────

def make_slide7(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "1ª Passagem — Frase 2",
                    "The team played its first international game in 1914.")

    headers = ["Palavra / trecho", "O que o aluno reconhece"]
    rows = [
        ["The team",      "o time — já apareceu antes"],
        ["played",        ("não sei a palavra, mas soa como passado", LIGHT_GRAY)],
        ["its first",     ("não sei", LIGHT_GRAY)],
        ["international", "internacional — claro"],
        ["game",          "jogo — já ouviu em outros contextos"],
        [("in 1914", ACCENT), "em 1914 — uma data"],
    ]
    col_widths = [Inches(3.5), Inches(9.0)]
    add_simple_table(slide, headers, rows,
                     x=Inches(0.4), y=Inches(1.25),
                     col_widths=col_widths, row_height=Inches(0.56))

    add_callout(slide,
                "O aluno entende: o time fez algo internacional em 1914. "
                "Provavelmente jogou algo. Faz sentido com o contexto anterior.",
                Inches(0.4), Inches(5.3), Inches(12.5), Inches(0.65),
                bg=RGBColor(0x1A, 0x3A, 0x2A), icon="✓")

    add_textbox(slide,
                "Nota: played its first é opaco para Tier 1, "
                "mas 1914 e international game são pontos de apoio suficientes.",
                Inches(0.4), Inches(6.1), Inches(12.5), Inches(0.55),
                font_size=15, color=LIGHT_GRAY, italic=True)

    return slide


# ── SLIDE 8 — Frase 3: Since then, Brazil has competed ──────────────────────

def make_slide8(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "1ª Passagem — Frase 3",
                    "Since then, Brazil has competed in every single World Cup.")

    headers = ["Palavra / trecho", "O que o aluno reconhece"]
    rows = [
        ["Since then",         ("não sei", LIGHT_GRAY)],
        [("Brazil", ACCENT),   "Brasil"],
        ["has competed",       "parece compete... competiu?"],
        ["in every single",    ("não sei every single, mas...", LIGHT_GRAY)],
        [("World Cup", ACCENT), "Copa do Mundo — claro"],
    ]
    col_widths = [Inches(3.5), Inches(9.0)]
    add_simple_table(slide, headers, rows,
                     x=Inches(0.4), y=Inches(1.25),
                     col_widths=col_widths, row_height=Inches(0.62))

    add_callout(slide,
                "O aluno entende: Brasil + Copa do Mundo + ideia de competir. "
                "Provavelmente 'o Brasil participou de toda Copa do Mundo'.",
                Inches(0.4), Inches(5.05), Inches(12.5), Inches(0.65),
                bg=RGBColor(0x1A, 0x3A, 0x2A), icon="✓")

    add_textbox(slide,
                "Frase mais opaca das quatro para Tier 1. Mesmo assim, "
                "Brazil e World Cup carregam o peso semântico.",
                Inches(0.4), Inches(5.88), Inches(12.5), Inches(0.55),
                font_size=15, color=LIGHT_GRAY, italic=True)

    return slide


# ── SLIDE 9 — Frase 4: No other team ────────────────────────────────────────

def make_slide9(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "1ª Passagem — Frase 4",
                    "No other team has done this.")

    headers = ["Palavra / trecho", "O que o aluno reconhece"]
    rows = [
        [("No", ACCENT),   "não, nenhum"],
        ["other",          "parece outro"],
        ["team",           "time — já apareceu duas vezes"],
        ["has done",       ("não sei", LIGHT_GRAY)],
        [("this", ACCENT), "isso — aponta para o que foi dito antes"],
    ]
    col_widths = [Inches(3.5), Inches(9.0)]
    add_simple_table(slide, headers, rows,
                     x=Inches(0.4), y=Inches(1.25),
                     col_widths=col_widths, row_height=Inches(0.62))

    add_callout(slide,
                "O aluno entende: nenhum outro time fez... alguma coisa. "
                "O this remete ao contexto anterior — participar de toda Copa. "
                "O aluno provavelmente capta a ligação.",
                Inches(0.4), Inches(5.05), Inches(12.5), Inches(0.72),
                bg=RGBColor(0x1A, 0x3A, 0x2A), icon="✓")

    add_textbox(slide, "O aluno sente satisfação — entendeu o argumento geral do parágrafo.",
                Inches(0.4), Inches(5.95), Inches(12.5), Inches(0.5),
                font_size=16, color=GREEN, bold=True)

    return slide


# ── SLIDE 10 — O que o aluno entendeu (resultado da 1ª passagem) ─────────────

def make_slide10(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "Resultado da 1ª passagem — o que o aluno entendeu")

    add_code_block(slide,
                   '"É sobre a Seleção Brasileira. Ela participou de todas as Copas do\n'
                   'Mundo desde 1914 e nenhum outro time fez isso."',
                   Inches(0.5), Inches(1.3), Inches(12.2), Inches(1.3), font_size=20)

    add_callout(slide,
                "Isso é suficiente para a 1ª passagem. "
                "Não é uma tradução frase a frase — é o sentido global do parágrafo.",
                Inches(0.4), Inches(2.85), Inches(12.5), Inches(0.65),
                bg=RGBColor(0x1A, 0x3A, 0x2A), icon="✓")

    add_textbox(slide, "O que a 1ª passagem NÃO faz (e não precisa fazer):",
                Inches(0.4), Inches(3.7), Inches(12.5), Inches(0.45),
                font_size=18, bold=True, color=ACCENT)

    headers = ["O aluno NÃO precisa saber", "Por que não importa agora"]
    rows = [
        ["O que é is called",     "O nome A Seleção já entrega o sentido"],
        ["O que é played its first", "1914 e international game são suficientes"],
        ["O que é since then",    "O contexto de sequência temporal já está claro"],
        ["O que é every single",  "World Cup é o ponto de apoio suficiente"],
        ["O que é has done this", "O this remete ao que veio antes — entendível pelo contexto"],
    ]
    col_widths = [Inches(4.0), Inches(8.7)]
    add_simple_table(slide, headers, rows,
                     x=Inches(0.4), y=Inches(4.2),
                     col_widths=col_widths, row_height=Inches(0.46))

    return slide


# ── SLIDE 11 — Por que o aluno não vai usar Google Translate ─────────────────

def make_slide11(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "Por que o aluno não vai usar o Google Translate")
    add_textbox(slide, "(se a instrução for dada do jeito certo)",
                Inches(0.4), Inches(1.15), Inches(12.5), Inches(0.4),
                font_size=17, color=LIGHT_GRAY, italic=True)

    # Coluna esquerda: O que NÃO funciona
    add_rect(slide, Inches(0.4), Inches(1.7), Inches(5.9), Inches(3.2),
             RGBColor(0x3A, 0x1A, 0x1A))
    add_textbox(slide, "✗  O que NÃO funciona:",
                Inches(0.55), Inches(1.82), Inches(5.6), Inches(0.45),
                font_size=18, bold=True, color=RED)
    add_code_block(slide,
                   '"Leia sem usar\no Google Translate."',
                   Inches(0.55), Inches(2.35), Inches(5.6), Inches(0.9), font_size=14)
    add_textbox(slide,
                "Só diz o que não pode fazer.\nNão diz o que fazer em vez disso.",
                Inches(0.55), Inches(3.35), Inches(5.6), Inches(0.85),
                font_size=16, color=LIGHT_GRAY, italic=True)

    # Coluna direita: O que funciona
    add_rect(slide, Inches(7.0), Inches(1.7), Inches(5.9), Inches(3.2),
             RGBColor(0x1A, 0x3A, 0x2A))
    add_textbox(slide, "✓  O que funciona:",
                Inches(7.15), Inches(1.82), Inches(5.6), Inches(0.45),
                font_size=18, bold=True, color=GREEN)
    add_code_block(slide,
                   '"Leia uma vez do início ao fim.\n'
                   'Quando terminar, escreva em uma\n'
                   'frase o que o texto está falando."',
                   Inches(7.15), Inches(2.35), Inches(5.6), Inches(1.25), font_size=13)
    add_textbox(slide,
                "Dá uma tarefa concreta que o aluno\nconsegue executar sem Google Translate.",
                Inches(7.15), Inches(3.7), Inches(5.6), Inches(0.7),
                font_size=16, color=LIGHT_GRAY, italic=True)

    add_callout(slide,
                "A pergunta muda de 'qual é a tradução de cada palavra?' "
                "para 'sobre o que é esse texto?' — e essa tem resposta com vocabulário limitado.",
                Inches(0.4), Inches(5.25), Inches(12.5), Inches(0.72),
                bg=RGBColor(0x1E, 0x2E, 0x4E), icon="→")

    return slide


# ── SLIDE 12 — Conexão entre 1ª e 2ª passagem ───────────────────────────────

def make_slide12(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "Conexão entre as duas passagens")

    add_textbox(slide,
                "A 1ª passagem cria o contexto que torna a 2ª passagem produtiva.",
                Inches(0.4), Inches(1.25), Inches(12.5), Inches(0.5),
                font_size=20, bold=True, color=ACCENT)

    # Caixa 1ª passagem
    add_rect(slide, Inches(0.4), Inches(2.0), Inches(5.8), Inches(1.8),
             RGBColor(0x1E, 0x2E, 0x4E))
    add_textbox(slide, "1ª passagem",
                Inches(0.55), Inches(2.1), Inches(5.5), Inches(0.45),
                font_size=18, bold=True, color=ACCENT)
    add_textbox(slide,
                "O aluno entende que o texto é sobre\na Seleção Brasileira e Copa do Mundo.",
                Inches(0.55), Inches(2.6), Inches(5.5), Inches(1.0),
                font_size=16, color=WHITE)

    # Seta
    add_textbox(slide, "→", Inches(6.5), Inches(2.7), Inches(0.8), Inches(0.6),
                font_size=38, bold=True, color=ACCENT)

    # Caixa 2ª passagem
    add_rect(slide, Inches(7.55), Inches(2.0), Inches(5.4), Inches(1.8),
             RGBColor(0x1A, 0x3A, 0x2A))
    add_textbox(slide, "2ª passagem",
                Inches(7.7), Inches(2.1), Inches(5.1), Inches(0.45),
                font_size=18, bold=True, color=GREEN)
    add_textbox(slide,
                '"Ache a palavra played no texto.\n'
                'Em qual frase ela está? Quem está\n'
                'fazendo a ação? O que o time fez?"',
                Inches(7.7), Inches(2.6), Inches(5.1), Inches(1.05),
                font_size=15, color=WHITE, italic=True)

    add_callout(slide,
                "Análise sintática com propósito e contexto — não análise no vazio.",
                Inches(0.4), Inches(4.1), Inches(12.5), Inches(0.65),
                bg=RGBColor(0x1E, 0x2E, 0x4E), icon="→")

    add_textbox(slide,
                "Somente depois que o aluno tem o modelo mental global do texto "
                "é que faz sentido introduzir análise frase a frase.",
                Inches(0.4), Inches(4.95), Inches(12.5), Inches(0.55),
                font_size=17, color=LIGHT_GRAY, italic=True)

    add_textbox(slide,
                "Na 2ª passagem: a análise sintática entra APENAS nas frases "
                "que contêm as palavras da lista do professor.",
                Inches(0.4), Inches(5.65), Inches(12.5), Inches(0.55),
                font_size=17, color=WHITE)

    return slide


# ── SLIDE 13 — Divisor: 2ª Passagem ─────────────────────────────────────────

def make_slide13(prs):
    return section_divider(prs,
        "2ª Passagem de Leitura",
        "Caçando as palavras-alvo dentro do texto")


# ── SLIDE 14 — O que mudou desde a 1ª passagem ──────────────────────────────

def make_slide14(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "O que mudou desde a 1ª passagem")

    add_textbox(slide,
                "Agora o aluno tem duas coisas que não tinha antes:",
                Inches(0.4), Inches(1.25), Inches(12.5), Inches(0.5),
                font_size=20, color=WHITE)

    # Item 1
    add_rect(slide, Inches(0.4), Inches(1.9), Inches(12.0), Inches(1.5),
             RGBColor(0x1E, 0x2E, 0x4E))
    add_textbox(slide, "① Um modelo mental do texto",
                Inches(0.6), Inches(2.0), Inches(11.5), Inches(0.45),
                font_size=20, bold=True, color=ACCENT)
    add_textbox(slide,
                "Sabe do que se trata, não vai encontrar nenhuma surpresa de conteúdo.",
                Inches(0.6), Inches(2.5), Inches(11.5), Inches(0.6),
                font_size=17, color=WHITE)

    # Item 2
    add_rect(slide, Inches(0.4), Inches(3.6), Inches(12.0), Inches(1.5),
             RGBColor(0x1A, 0x3A, 0x2A))
    add_textbox(slide, "② A lista de 10 palavras-alvo do professor",
                Inches(0.6), Inches(3.7), Inches(11.5), Inches(0.45),
                font_size=20, bold=True, color=GREEN)
    add_textbox(slide,
                "Sabe exatamente onde concentrar a atenção.",
                Inches(0.6), Inches(4.2), Inches(11.5), Inches(0.6),
                font_size=17, color=WHITE)

    add_callout(slide,
                "O aluno não relê o texto inteiro com atenção igual. "
                "Ele caça as palavras da lista e examina cada uma dentro da sua frase.",
                Inches(0.4), Inches(5.35), Inches(12.5), Inches(0.72),
                bg=MID_GRAY, icon="→")

    return slide


# ── SLIDE 15 — A instrução da 2ª passagem ───────────────────────────────────

def make_slide15(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "A instrução correta para o aluno — 2ª passagem")

    add_code_block(slide,
                   '"Agora você vai reler o texto com a lista de palavras na mão.\n'
                   'Para cada palavra da lista:\n'
                   '  (1) ache a frase do texto onde ela aparece;\n'
                   '  (2) identifique o verbo da frase e o tempo verbal;\n'
                   '  (3) identifique quem está fazendo a ação;\n'
                   '  (4) identifique o que ou quem sofre a ação.\n'
                   'Só depois consulte o dicionário para essa palavra."',
                   Inches(0.4), Inches(1.4), Inches(12.5), Inches(2.8), font_size=17)

    add_callout(slide,
                "A análise verbo → sujeito → objeto é feita frase a frase, "
                "SOMENTE nas frases das palavras-alvo.",
                Inches(0.4), Inches(4.45), Inches(12.5), Inches(0.65),
                bg=RGBColor(0x1A, 0x3A, 0x2A), icon="✓")

    add_textbox(slide,
                "Por que verbo primeiro?",
                Inches(0.4), Inches(5.25), Inches(12.5), Inches(0.4),
                font_size=17, bold=True, color=ACCENT)

    items = [
        "O verbo diz o que está acontecendo — sem ele, a frase não existe",
        "O tempo verbal diz quando — presente, passado, futuro",
        "Sujeito e objeto só fazem sentido a partir do verbo",
    ]
    for i, item in enumerate(items):
        add_textbox(slide, f"• {item}",
                    Inches(0.6), Inches(5.75) + Inches(0.38) * i,
                    Inches(12.1), Inches(0.36),
                    font_size=16, color=WHITE)

    return slide


# ── SLIDE 16 — Palavra 1: won ────────────────────────────────────────────────

def make_slide16(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "2ª Passagem — Palavra 1: won",
                    "Brazil won its first World Cup in 1958 in Sweden.")

    headers = ["Pergunta", "Resposta do aluno"]
    rows = [
        ["Qual é o verbo?",        ("won", ACCENT)],
        ["Que tempo verbal?",      "Parece passado — a data 1958 confirma"],
        ["Quem fez a ação?",       "Brazil"],
        ["O que Brazil fez?",      "its first World Cup — ganhou a Copa"],
    ]
    col_widths = [Inches(4.0), Inches(8.6)]
    add_simple_table(slide, headers, rows,
                     x=Inches(0.4), y=Inches(1.25),
                     col_widths=col_widths, row_height=Inches(0.68))

    add_callout(slide,
                'Antes do dicionário: "O Brasil fez algo com a Copa do Mundo em 1958. '
                'Pelo contexto, ganhou."',
                Inches(0.4), Inches(4.15), Inches(12.5), Inches(0.65),
                bg=RGBColor(0x1A, 0x3A, 0x2A), icon="→")

    add_code_block(slide,
                   'Frente: "Brazil _______ its first World Cup in 1958."\n'
                   'Verso:  won — ganhou\n'
                   '        Minha frase: "My team won the game yesterday."',
                   Inches(0.4), Inches(5.0), Inches(12.5), Inches(1.35), font_size=15)

    return slide


# ── SLIDE 17 — Palavra 2: lifted ────────────────────────────────────────────

def make_slide17(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "2ª Passagem — Palavra 2: lifted",
                    "Brazil defeated Sweden in the final and lifted the trophy.")

    headers = ["Pergunta", "Resposta do aluno"]
    rows = [
        ["Qual é o verbo?",     ("Há dois: defeated e lifted — separados por and", ACCENT)],
        ["Que tempo verbal?",   "Passado — igual a won"],
        ["Quem fez as ações?",  "Brazil"],
        ["O que Brazil fez primeiro?", "defeated Sweden"],
        ["O que Brazil fez depois?",   "lifted the trophy"],
    ]
    col_widths = [Inches(3.8), Inches(8.8)]
    add_simple_table(slide, headers, rows,
                     x=Inches(0.4), y=Inches(1.25),
                     col_widths=col_widths, row_height=Inches(0.62))

    add_callout(slide,
                "Nota: Quando você vê 'and' entre dois verbos, o sujeito fez as duas coisas.",
                Inches(0.4), Inches(4.65), Inches(12.5), Inches(0.6),
                bg=RGBColor(0x1E, 0x2E, 0x4E), icon="ℹ")

    add_code_block(slide,
                   'Frente: "Brazil defeated Sweden and _______ the trophy."\n'
                   'Verso:  lifted — levantou, ergueu\n'
                   '        Colocação: lift the trophy\n'
                   '        Minha frase: "The captain lifted the trophy after the final."',
                   Inches(0.4), Inches(5.4), Inches(12.5), Inches(1.4), font_size=14)

    return slide


# ── SLIDE 18 — Palavra 3: defeated ──────────────────────────────────────────

def make_slide18(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "2ª Passagem — Palavra 3: defeated",
                    "Brazil defeated Sweden in the final and lifted the trophy.")

    headers = ["Pergunta", "Resposta do aluno"]
    rows = [
        ["Qual é o verbo desta vez?", ("defeated", ACCENT)],
        ["Que tempo verbal?",          "Passado"],
        ["Quem fez a ação?",           "Brazil"],
        ["Quem sofreu a ação?",        "Sweden"],
    ]
    col_widths = [Inches(3.8), Inches(8.8)]
    add_simple_table(slide, headers, rows,
                     x=Inches(0.4), y=Inches(1.25),
                     col_widths=col_widths, row_height=Inches(0.68))

    add_callout(slide,
                'Antes do dicionário: "O Brasil fez algo com a Suécia na final. '
                'Pelo contexto, ganhou dela — derrotou."',
                Inches(0.4), Inches(4.15), Inches(12.5), Inches(0.65),
                bg=RGBColor(0x1A, 0x3A, 0x2A), icon="→")

    add_code_block(slide,
                   'Frente: "Brazil _______ Sweden in the final."\n'
                   'Verso:  defeated — derrotou\n'
                   '        Colocação: defeat a team / defeat an opponent\n'
                   '        Minha frase: "Brazil defeated Argentina in the semifinals."',
                   Inches(0.4), Inches(5.0), Inches(12.5), Inches(1.4), font_size=14)

    return slide


# ── SLIDE 19 — Palavra 4: competed ──────────────────────────────────────────

def make_slide19(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "2ª Passagem — Palavra 4: competed",
                    "Since then, Brazil has competed in every single World Cup.")

    headers = ["Pergunta", "Resposta do aluno"]
    rows = [
        ["Qual é o verbo?",       ("has competed — parece compete numa forma diferente", ACCENT)],
        ["Que tempo verbal?",      ("Não é passado simples como os outros...", LIGHT_GRAY)],
        ["Quem fez a ação?",       "Brazil"],
        ["Em que Brazil competiu?","every single World Cup"],
    ]
    col_widths = [Inches(3.8), Inches(8.8)]
    add_simple_table(slide, headers, rows,
                     x=Inches(0.4), y=Inches(1.25),
                     col_widths=col_widths, row_height=Inches(0.68))

    add_callout(slide,
                "Oportunidade: has competed é present perfect — mesmo tempo de E01. "
                "'O Brasil começou a competir no passado e ainda compete hoje.'",
                Inches(0.4), Inches(4.15), Inches(12.5), Inches(0.72),
                bg=RGBColor(0x1E, 0x2E, 0x4E), icon="ℹ")

    add_code_block(slide,
                   'Frente: "Brazil has _______ in every single World Cup."\n'
                   'Verso:  competed — competiu, participou\n'
                   '        Colocação: compete in the World Cup\n'
                   '        Minha frase: "My university competed in the national championship."',
                   Inches(0.4), Inches(5.1), Inches(12.5), Inches(1.4), font_size=14)

    return slide


# ── SLIDE 20 — Palavra 5: squad ─────────────────────────────────────────────

def make_slide20(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "2ª Passagem — Palavra 5: squad",
                    "Many experts consider the 1970 squad the best team ever.")

    headers = ["Pergunta", "Resposta do aluno"]
    rows = [
        ["Qual é o verbo?",         ("consider", ACCENT)],
        ["Que tempo verbal?",        "Presente simples"],
        ["Quem fez a ação?",         "Many experts — muitos especialistas"],
        ["O que os especialistas\nconsideram?", "the 1970 squad the best team ever"],
    ]
    col_widths = [Inches(3.8), Inches(8.8)]
    add_simple_table(slide, headers, rows,
                     x=Inches(0.4), y=Inches(1.25),
                     col_widths=col_widths, row_height=Inches(0.68))

    add_callout(slide,
                'Antes do dicionário: "Especialistas consideram o ____ de 1970 '
                'o melhor time de todos os tempos. Squad deve ser elenco."',
                Inches(0.4), Inches(4.15), Inches(12.5), Inches(0.65),
                bg=RGBColor(0x1A, 0x3A, 0x2A), icon="→")

    add_code_block(slide,
                   'Frente: "Many experts consider the 1970 _______ the best team ever."\n'
                   'Verso:  squad — elenco, seleção\n'
                   '        Minha frase: "The coach announced the squad for the World Cup."',
                   Inches(0.4), Inches(5.0), Inches(12.5), Inches(1.3), font_size=14)

    return slide


# ── SLIDE 21 — Palavra 6: title ─────────────────────────────────────────────

def make_slide21(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "2ª Passagem — Palavra 6: title",
                    "In total, Brazil has won five World Cup titles.")

    headers = ["Pergunta", "Resposta do aluno"]
    rows = [
        ["Qual é o verbo?",     ("has won", ACCENT)],
        ["Que tempo verbal?",    "Present perfect novamente"],
        ["Quem fez a ação?",     "Brazil"],
        ["O que Brazil ganhou?", "five World Cup titles"],
    ]
    col_widths = [Inches(3.8), Inches(8.8)]
    add_simple_table(slide, headers, rows,
                     x=Inches(0.4), y=Inches(1.25),
                     col_widths=col_widths, row_height=Inches(0.68))

    add_callout(slide,
                "title é cognato parcial — o aluno pode reconhecer. "
                "O valor está na colocação: win the title / World Cup title.",
                Inches(0.4), Inches(4.15), Inches(12.5), Inches(0.65),
                bg=RGBColor(0x1E, 0x2E, 0x4E), icon="ℹ")

    add_code_block(slide,
                   'Frente: "Brazil has won five World Cup _______."  \n'
                   'Verso:  title — título\n'
                   '        Colocação: win a title / World Cup title\n'
                   '        Minha frase: "Brazil wants to win a sixth title in 2026."',
                   Inches(0.4), Inches(5.0), Inches(12.5), Inches(1.35), font_size=14)

    return slide


# ── SLIDE 22 — Palavra 7: rival ─────────────────────────────────────────────

def make_slide22(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "2ª Passagem — Palavra 7: rival",
                    "Brazil's greatest rival is Argentina.")

    headers = ["Pergunta", "Resposta do aluno"]
    rows = [
        ["Qual é o verbo?",        ("is", ACCENT)],
        ["Que tempo verbal?",       "Presente simples"],
        ["Quem é o sujeito?",       "Brazil's greatest rival"],
        ["O que o verbo diz?",      "is Argentina — iguala os dois"],
    ]
    col_widths = [Inches(3.8), Inches(8.8)]
    add_simple_table(slide, headers, rows,
                     x=Inches(0.4), y=Inches(1.25),
                     col_widths=col_widths, row_height=Inches(0.68))

    add_callout(slide,
                "rival é cognato direto — o aluno reconhece na hora. "
                "O valor do cartão está na colocação greatest rival.",
                Inches(0.4), Inches(4.15), Inches(12.5), Inches(0.65),
                bg=RGBColor(0x1E, 0x2E, 0x4E), icon="ℹ")

    add_code_block(slide,
                   'Frente: "Brazil\'s _______ rival is Argentina."\n'
                   'Verso:  greatest rival — maior rival\n'
                   '        Colocação: greatest rival / fierce rival\n'
                   '        Minha frase: "Real Madrid\'s greatest rival is Barcelona."',
                   Inches(0.4), Inches(5.0), Inches(12.5), Inches(1.35), font_size=14)

    return slide


# ── SLIDE 23 — Palavra 8: fiercely ──────────────────────────────────────────

def make_slide23(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "2ª Passagem — Palavra 8: fiercely",
                    "The two teams compete fiercely.")

    headers = ["Pergunta", "Resposta do aluno"]
    rows = [
        ["Qual é o verbo?",   ("compete", ACCENT)],
        ["Que tempo verbal?",  "Presente simples"],
        ["Quem faz a ação?",   "The two teams"],
        ["Como eles competem?", ("fiercely", ACCENT)],
    ]
    col_widths = [Inches(3.8), Inches(8.8)]
    add_simple_table(slide, headers, rows,
                     x=Inches(0.4), y=Inches(1.25),
                     col_widths=col_widths, row_height=Inches(0.68))

    add_callout(slide,
                "fiercely é advérbio — modifica o verbo. O aluno não precisa saber disso. "
                "A pergunta 'como eles competem?' aponta naturalmente para ele.",
                Inches(0.4), Inches(4.15), Inches(12.5), Inches(0.65),
                bg=RGBColor(0x1E, 0x2E, 0x4E), icon="ℹ")

    add_code_block(slide,
                   'Frente: "The two teams compete _______."  \n'
                   'Verso:  fiercely — ferozmente, com intensidade\n'
                   '        Colocação: compete fiercely\n'
                   '        Minha frase: "The students competed fiercely for the scholarship."',
                   Inches(0.4), Inches(5.0), Inches(12.5), Inches(1.35), font_size=14)

    return slide


# ── SLIDE 24 — Palavra 9: solid ─────────────────────────────────────────────

def make_slide24(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "2ª Passagem — Palavra 9: solid",
                    "The team has talented players and a solid squad.")

    headers = ["Pergunta", "Resposta do aluno"]
    rows = [
        ["Qual é o verbo?",   ("has", ACCENT)],
        ["Que tempo verbal?",  "Presente simples"],
        ["Quem faz a ação?",   "The team"],
        ["O que o time tem?",  "talented players e a solid squad"],
    ]
    col_widths = [Inches(3.8), Inches(8.8)]
    add_simple_table(slide, headers, rows,
                     x=Inches(0.4), y=Inches(1.25),
                     col_widths=col_widths, row_height=Inches(0.68))

    add_callout(slide,
                "solid parece sólido — o aluno pode inferir: forte, consistente. "
                "Em contexto esportivo: bom coletivo, não depende de um só jogador.",
                Inches(0.4), Inches(4.15), Inches(12.5), Inches(0.65),
                bg=RGBColor(0x1A, 0x3A, 0x2A), icon="→")

    add_code_block(slide,
                   'Frente: "The team has talented players and a _______ squad."\n'
                   'Verso:  solid — sólido, forte, consistente\n'
                   '        Minha frase: "Brazil has a solid squad for the 2026 World Cup."',
                   Inches(0.4), Inches(5.0), Inches(12.5), Inches(1.3), font_size=14)

    return slide


# ── SLIDE 25 — Palavra 10: fans ─────────────────────────────────────────────

def make_slide25(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "2ª Passagem — Palavra 10: fans",
                    "Fans around the world hope Brazil will win a sixth title.")

    headers = ["Pergunta", "Resposta do aluno"]
    rows = [
        ["Qual é o verbo?",       ("hope", ACCENT)],
        ["Que tempo verbal?",      "Presente simples"],
        ["Quem faz a ação?",       ("Fans", ACCENT)],
        ["O que os fans esperam?", "Brazil will win a sixth title"],
    ]
    col_widths = [Inches(3.8), Inches(8.8)]
    add_simple_table(slide, headers, rows,
                     x=Inches(0.4), y=Inches(1.25),
                     col_widths=col_widths, row_height=Inches(0.68))

    add_callout(slide,
                "Fans devem ser torcedores — o aluno já ouviu essa palavra.",
                Inches(0.4), Inches(4.15), Inches(12.5), Inches(0.6),
                bg=RGBColor(0x1A, 0x3A, 0x2A), icon="→")

    add_code_block(slide,
                   'Frente: "_______ around the world hope Brazil will win a sixth title."\n'
                   'Verso:  fans — torcedores, fãs\n'
                   '        Colocação: fans around the world / football fans\n'
                   '        Minha frase: "Fans around the world watched the final on TV."',
                   Inches(0.4), Inches(4.95), Inches(12.5), Inches(1.4), font_size=14)

    return slide


# ── SLIDE 26 — O que a 2ª passagem entregou ─────────────────────────────────

def make_slide26(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "O que a 2ª passagem entregou")

    items = [
        ("①", "Localizou cada palavra dentro de uma frase real",
         "Não consultou o dicionário no vazio"),
        ("②", "Entendeu o papel de cada palavra na frase",
         "O dicionário confirmou — não revelou do zero"),
        ("③", "Criou 10 cartões com frente cloze e frase própria",
         "Pronto para a fase de revisão"),
        ("④", "Praticou as 3 perguntas de leitura de forma repetida",
         "Verbo → Quem → O quê — sem chamar isso de gramática"),
    ]

    for i, (num, title, desc) in enumerate(items):
        y = Inches(1.5) + Inches(1.3) * i
        add_rect(slide, Inches(0.4), y, Inches(12.5), Inches(1.1),
                 RGBColor(0x26, 0x26, 0x3A))
        add_textbox(slide, num,
                    Inches(0.55), y + Inches(0.12), Inches(0.55), Inches(0.8),
                    font_size=28, bold=True, color=ACCENT)
        add_textbox(slide, title,
                    Inches(1.25), y + Inches(0.1), Inches(11.4), Inches(0.45),
                    font_size=19, bold=True, color=WHITE)
        add_textbox(slide, desc,
                    Inches(1.25), y + Inches(0.58), Inches(11.4), Inches(0.38),
                    font_size=16, color=LIGHT_GRAY, italic=True)

    return slide


# ── SLIDE 27 — Resumo: as duas passagens ────────────────────────────────────

def make_slide27(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "Resumo: as duas passagens de leitura — Tier 1")

    headers = ["", "1ª Passagem", "2ª Passagem"]
    rows = [
        [("Objetivo", ACCENT),
         "Entender do que o texto trata",
         "Examinar as palavras-alvo dentro de suas frases"],
        [("Suporte", ACCENT),
         ("Nenhum", GREEN),
         "Lista do professor + dicionário palavra a palavra"],
        [("Foco", ACCENT),
         "Texto inteiro",
         "Somente as frases das 10 palavras-alvo"],
        [("Ferramenta", ACCENT),
         "Cognatos + contexto",
         ("Verbo → Sujeito → Objeto → Dicionário", ACCENT)],
        [("Produto", ACCENT),
         "Modelo mental do texto",
         ("10 cartões prontos para revisão", GREEN)],
        [("Tempo estimado", ACCENT),
         "5–8 minutos",
         "20–30 minutos"],
    ]
    col_widths = [Inches(2.2), Inches(4.9), Inches(5.8)]
    add_simple_table(slide, headers, rows,
                     x=Inches(0.4), y=Inches(1.15),
                     col_widths=col_widths, row_height=Inches(0.72))

    return slide


# ── SLIDE 28 — Instrução consolidada ────────────────────────────────────────

def make_slide28(prs):
    slide = blank_slide(prs)
    slide_title_bar(slide, "Instrução consolidada para o aluno Tier 1")

    # 1ª passagem
    add_textbox(slide, "1ª passagem:",
                Inches(0.4), Inches(1.3), Inches(12.5), Inches(0.42),
                font_size=18, bold=True, color=ACCENT)
    add_code_block(slide,
                   '"Leia o texto uma vez do início ao fim. Não pare nas palavras que não\n'
                   'conhece. Quando terminar, escreva em uma frase em português o que o texto\n'
                   'está falando."',
                   Inches(0.4), Inches(1.8), Inches(12.5), Inches(1.2), font_size=15)

    # 2ª passagem
    add_textbox(slide, "2ª passagem:",
                Inches(0.4), Inches(3.15), Inches(12.5), Inches(0.42),
                font_size=18, bold=True, color=GREEN)
    add_code_block(slide,
                   '"Agora você vai reler o texto com a lista de palavras na mão.\n'
                   'Para cada palavra: (1) ache a frase; (2) identifique o verbo;\n'
                   '(3) identifique quem faz a ação; (4) identifique o que ou quem sofre a ação.\n'
                   'Só depois consulte o dicionário."',
                   Inches(0.4), Inches(3.65), Inches(12.5), Inches(1.35), font_size=15)

    # Criação do cartão
    add_textbox(slide, "Criação do cartão:",
                Inches(0.4), Inches(5.15), Inches(12.5), Inches(0.42),
                font_size=18, bold=True, color=BLUE)
    add_code_block(slide,
                   '"Copie a frase do texto e retire a palavra (deixe um blank).\n'
                   'No verso: escreva a tradução + uma frase sua trocando 1 palavra por algo da sua vida."',
                   Inches(0.4), Inches(5.65), Inches(12.5), Inches(1.05), font_size=15)

    return slide


# ── Build & save ─────────────────────────────────────────────────────────────

def build():
    prs = new_prs()

    make_slide1(prs)     # Capa
    make_slide2(prs)     # Visão geral
    make_slide3(prs)     # O problema
    make_slide4(prs)     # Critério de sucesso
    make_slide5(prs)     # Texto de referência
    # Divisor 1ª passagem
    section_divider(prs, "1ª Passagem de Leitura",
                    "Leitura do início ao fim — sem parar, sem dicionário")
    make_slide6(prs)     # Frase 1
    make_slide7(prs)     # Frase 2
    make_slide8(prs)     # Frase 3
    make_slide9(prs)     # Frase 4
    make_slide10(prs)    # Resultado da 1ª passagem
    make_slide11(prs)    # Por que não usar Google Translate
    make_slide12(prs)    # Conexão 1ª → 2ª passagem
    make_slide13(prs)    # Divisor 2ª passagem
    make_slide14(prs)    # O que mudou
    make_slide15(prs)    # Instrução da 2ª passagem
    make_slide16(prs)    # Palavra 1: won
    make_slide17(prs)    # Palavra 2: lifted
    make_slide18(prs)    # Palavra 3: defeated
    make_slide19(prs)    # Palavra 4: competed
    make_slide20(prs)    # Palavra 5: squad
    make_slide21(prs)    # Palavra 6: title
    make_slide22(prs)    # Palavra 7: rival
    make_slide23(prs)    # Palavra 8: fiercely
    make_slide24(prs)    # Palavra 9: solid
    make_slide25(prs)    # Palavra 10: fans
    make_slide26(prs)    # O que a 2ª passagem entregou
    make_slide27(prs)    # Resumo comparativo
    make_slide28(prs)    # Instrução consolidada

    out = "E02_slides_Tier1_Duas_Passagem_v1.0.pptx"
    prs.save(out)
    print(f"Salvo: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
